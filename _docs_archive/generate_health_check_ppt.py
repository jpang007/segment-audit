#!/usr/bin/env python3
"""
Technical Health Check PowerPoint Generator
Takes Gem JSON output and generates a PowerPoint presentation with charts
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import json
import sys
from pathlib import Path
from datetime import datetime


class HealthCheckPPTGenerator:
    """Generate Technical Health Check PowerPoint from Gem JSON output"""

    # Segment brand colors
    SEGMENT_GREEN = RGBColor(82, 204, 138)
    SEGMENT_BLUE = RGBColor(82, 138, 204)
    SEGMENT_ORANGE = RGBColor(255, 138, 82)
    SEGMENT_RED = RGBColor(255, 82, 82)
    SEGMENT_PURPLE = RGBColor(138, 82, 204)
    SEGMENT_YELLOW = RGBColor(255, 204, 82)

    CHART_COLORS = [
        RGBColor(82, 138, 204),   # Blue
        RGBColor(82, 204, 138),   # Green
        RGBColor(255, 138, 82),   # Orange
        RGBColor(138, 82, 204),   # Purple
        RGBColor(255, 204, 82),   # Yellow
        RGBColor(255, 82, 82),    # Red
    ]

    def __init__(self, gem_json_path: str, customer_name: str = "Customer"):
        """Initialize with Gem JSON output"""
        with open(gem_json_path, 'r') as f:
            self.data = json.load(f)

        self.customer_name = customer_name
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)  # 16:9 aspect ratio

    def generate(self, output_path: str = None):
        """Generate complete PowerPoint"""
        if output_path is None:
            timestamp = datetime.now().strftime("%Y%m%d")
            output_path = f"Health_Check_{self.customer_name}_{timestamp}.pptx"

        print(f"🎨 Generating Health Check PowerPoint for {self.customer_name}...")

        # Generate all slides
        self._add_title_slide()
        self._add_limits_slide()
        self._add_events_slide()
        self._add_syntax_slide()
        self._add_source_variety_slide()
        self._add_source_volume_slide()
        self._add_destination_variety_slide()
        self._add_connections_slide()
        self._add_team_slide()
        self._add_summary_slide()
        self._add_conclusions_slide()
        self._add_thank_you_slide()

        # Save
        self.prs.save(output_path)
        print(f"✅ PowerPoint saved: {output_path}")
        return output_path

    def _add_title_slide(self):
        """Slide 1: Title slide"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        title_frame = title.text_frame
        title_frame.text = f"{self.customer_name}\nTechnical Health Check"
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(33, 33, 33)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Date
        date = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.5))
        date_frame = date.text_frame
        date_frame.text = f"Prepared {datetime.now().strftime('%B %d, %Y')}"
        date_frame.paragraphs[0].font.size = Pt(18)
        date_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
        date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_limits_slide(self):
        """Slide 2: Data Limits"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        data = self.data.get('slide_2_limits', {})

        # Title
        self._add_slide_title(slide, data.get('title', 'Data: Limits'))

        # MTU Usage - Big Number
        mtu = data.get('key_metrics', {}).get('mtu_usage', {})
        percent = mtu.get('percent_consumed', 'N/A')

        percent_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(3), Inches(1.5))
        pf = percent_box.text_frame
        pf.text = percent
        pf.paragraphs[0].font.size = Pt(72)
        pf.paragraphs[0].font.bold = True
        pf.paragraphs[0].font.color.rgb = self.SEGMENT_GREEN
        pf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # MTU Details
        current = mtu.get('current', 0)
        quota = mtu.get('quota', 0)
        details = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(3), Inches(0.5))
        df = details.text_frame
        df.text = f"{current:,} / {quota:,} MTUs"
        df.paragraphs[0].font.size = Pt(16)
        df.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Insight
        insight = data.get('insight', '')
        insight_box = slide.shapes.add_textbox(Inches(4.5), Inches(2), Inches(4.5), Inches(1.5))
        if_frame = insight_box.text_frame
        if_frame.text = insight
        if_frame.word_wrap = True
        if_frame.paragraphs[0].font.size = Pt(14)
        if_frame.paragraphs[0].font.color.rgb = RGBColor(66, 66, 66)

    def _add_events_slide(self):
        """Slide 3: Events & Props"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        data = self.data.get('slide_3_events', {})
        self._add_slide_title(slide, data.get('title', 'Data: Events & Props'))

        # Top Events Chart (horizontal bar)
        top_events = data.get('charts', {}).get('top_events_by_volume', [])[:10]
        if top_events:
            chart_data = CategoryChartData()
            chart_data.categories = [e['event_name'] for e in top_events]
            chart_data.add_series('Volume', [e['volume'] for e in top_events])

            x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(5), Inches(3.5)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
            ).chart

            chart.has_legend = False
            chart.has_title = True
            chart.chart_title.text_frame.text = "Top Events by Volume"

            # Color bars
            series = chart.series[0]
            for i, point in enumerate(series.points):
                fill = point.format.fill
                fill.solid()
                fill.fore_color.rgb = self.CHART_COLORS[i % len(self.CHART_COLORS)]

        # Unique Events Chart
        unique_events = data.get('charts', {}).get('unique_events_by_source', [])[:8]
        if unique_events:
            chart_data = CategoryChartData()
            chart_data.categories = [e['source'] for e in unique_events]
            chart_data.add_series('Events', [e['event_count'] for e in unique_events])

            x, y, cx, cy = Inches(5.5), Inches(1.5), Inches(4), Inches(2.5)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart

            chart.has_legend = False
            chart.has_title = True
            chart.chart_title.text_frame.text = "Unique Events by Source"

            # Color bars
            series = chart.series[0]
            for i, point in enumerate(series.points):
                fill = point.format.fill
                fill.solid()
                fill.fore_color.rgb = self.SEGMENT_BLUE

        # Insight
        insight = data.get('insight', '')
        self._add_insight_box(slide, insight)

    def _add_syntax_slide(self):
        """Slide 4: Syntax Validation"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        data = self.data.get('slide_4_syntax', {})
        self._add_slide_title(slide, data.get('title', 'Data: Syntax'))

        validation = data.get('validation_results', {})

        # Health Score Badge
        health = validation.get('health_score', 'Unknown')
        color_map = {
            'Good': self.SEGMENT_GREEN,
            'Fair': self.SEGMENT_YELLOW,
            'Poor': self.SEGMENT_RED
        }

        badge = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(2), Inches(1))
        bf = badge.text_frame
        bf.text = f"Health: {health}"
        bf.paragraphs[0].font.size = Pt(24)
        bf.paragraphs[0].font.bold = True
        bf.paragraphs[0].font.color.rgb = color_map.get(health, RGBColor(102, 102, 102))
        bf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Issues List
        issues = validation.get('issues_found', [])
        y_position = 1.5

        for issue in issues[:3]:  # Top 3 issues
            # Issue type and count
            issue_box = slide.shapes.add_textbox(Inches(3.5), Inches(y_position), Inches(5.5), Inches(0.4))
            ib_frame = issue_box.text_frame
            ib_frame.text = f"❌ {issue['type']}: {issue['count']} events"
            ib_frame.paragraphs[0].font.size = Pt(14)
            ib_frame.paragraphs[0].font.bold = True

            # Examples
            examples_text = "Examples: " + ", ".join(issue['examples'][:3])
            ex_box = slide.shapes.add_textbox(Inches(3.5), Inches(y_position + 0.3), Inches(5.5), Inches(0.3))
            ex_frame = ex_box.text_frame
            ex_frame.text = examples_text
            ex_frame.paragraphs[0].font.size = Pt(11)
            ex_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)

            # Recommendation
            rec_box = slide.shapes.add_textbox(Inches(3.5), Inches(y_position + 0.6), Inches(5.5), Inches(0.4))
            rec_frame = rec_box.text_frame
            rec_frame.text = f"→ {issue['recommendation']}"
            rec_frame.paragraphs[0].font.size = Pt(11)
            rec_frame.paragraphs[0].font.italic = True
            rec_frame.paragraphs[0].font.color.rgb = self.SEGMENT_BLUE
            rec_frame.word_wrap = True

            y_position += 1.1

        # Insight
        insight = data.get('insight', '')
        self._add_insight_box(slide, insight)

    def _add_source_variety_slide(self):
        """Slide 5: Source Variety"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        data = self.data.get('slide_5_source_variety', {})
        self._add_slide_title(slide, data.get('title', 'Stack: Source Variety'))

        # Sources by Library Chart
        libraries = data.get('charts', {}).get('sources_by_library', [])
        if libraries:
            chart_data = CategoryChartData()
            chart_data.categories = [lib['library'] for lib in libraries]
            chart_data.add_series('Sources', [lib['count'] for lib in libraries])

            x, y, cx, cy = Inches(1), Inches(1.5), Inches(5), Inches(3)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart

            chart.has_legend = False
            chart.has_title = True
            chart.chart_title.text_frame.text = "Active Sources by Library"

            # Color bars
            series = chart.series[0]
            for i, point in enumerate(series.points):
                fill = point.format.fill
                fill.solid()
                fill.fore_color.rgb = self.CHART_COLORS[i % len(self.CHART_COLORS)]

        # Total Active Sources
        total = data.get('charts', {}).get('total_active_sources', 0)
        categories = data.get('charts', {}).get('total_categories', 0)

        stats = slide.shapes.add_textbox(Inches(6.5), Inches(2), Inches(3), Inches(1.5))
        sf = stats.text_frame
        sf.text = f"{total}\nActive Sources\n\n{categories}\nCategories"
        sf.paragraphs[0].font.size = Pt(32)
        sf.paragraphs[0].font.bold = True
        sf.paragraphs[0].font.color.rgb = self.SEGMENT_BLUE
        sf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Insight
        insight = data.get('insight', '')
        self._add_insight_box(slide, insight)

    def _add_source_volume_slide(self):
        """Slide 6: Source Volume"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        data = self.data.get('slide_6_source_volume', {})
        self._add_slide_title(slide, data.get('title', 'Stack: Source Volume'))

        # Top Sources Chart
        top_sources = data.get('charts', {}).get('top_sources_by_events', [])[:8]
        if top_sources:
            chart_data = CategoryChartData()
            chart_data.categories = [s['source'] for s in top_sources]
            chart_data.add_series('Events', [s['event_count'] for s in top_sources])

            x, y, cx, cy = Inches(1), Inches(1.5), Inches(7), Inches(3.5)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
            ).chart

            chart.has_legend = False
            chart.has_title = False

            # Color bars
            series = chart.series[0]
            for i, point in enumerate(series.points):
                fill = point.format.fill
                fill.solid()
                fill.fore_color.rgb = self.SEGMENT_GREEN

        # Insight
        insight = data.get('insight', '')
        self._add_insight_box(slide, insight)

    def _add_destination_variety_slide(self):
        """Slide 7: Destination Variety"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        data = self.data.get('slide_7_destination_variety', {})
        self._add_slide_title(slide, data.get('title', 'Stack: Destination Variety'))

        # Destinations by Category Chart
        categories = data.get('charts', {}).get('destinations_by_category', [])
        if categories:
            chart_data = CategoryChartData()
            chart_data.categories = [cat['category'] for cat in categories]
            chart_data.add_series('Destinations', [cat['count'] for cat in categories])

            x, y, cx, cy = Inches(1), Inches(1.5), Inches(6), Inches(3)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart

            chart.has_legend = False
            chart.has_title = False

            # Color bars
            series = chart.series[0]
            for i, point in enumerate(series.points):
                fill = point.format.fill
                fill.solid()
                fill.fore_color.rgb = self.CHART_COLORS[i % len(self.CHART_COLORS)]

        # Total Stats
        total = data.get('charts', {}).get('total_active_destinations', 0)
        total_categories = data.get('charts', {}).get('total_categories', 0)

        stats = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(2.5), Inches(1.5))
        sf = stats.text_frame
        sf.text = f"{total}\nActive\nDestinations"
        sf.paragraphs[0].font.size = Pt(28)
        sf.paragraphs[0].font.bold = True
        sf.paragraphs[0].font.color.rgb = self.SEGMENT_PURPLE
        sf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Insight
        insight = data.get('insight', '')
        self._add_insight_box(slide, insight)

    def _add_connections_slide(self):
        """Slide 8: Connections"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        data = self.data.get('slide_8_connections', {})
        self._add_slide_title(slide, data.get('title', 'Stack: Source↔Destination Connections'))

        # Total Connections - Big Number
        total = data.get('charts', {}).get('total_connections', 0)

        total_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(2), Inches(1.5))
        tf = total_box.text_frame
        tf.text = f"{total}\nConnections"
        tf.paragraphs[0].font.size = Pt(48)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self.SEGMENT_ORANGE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Top Destinations
        top_dests = data.get('charts', {}).get('top_destinations_by_sources', [])[:5]
        if top_dests:
            chart_data = CategoryChartData()
            chart_data.categories = [d['destination'] for d in top_dests]
            chart_data.add_series('Sources', [d['connected_sources'] for d in top_dests])

            x, y, cx, cy = Inches(4.5), Inches(1.5), Inches(5), Inches(3)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
            ).chart

            chart.has_legend = False
            chart.has_title = True
            chart.chart_title.text_frame.text = "Top Destinations by Source Count"

            # Color bars
            series = chart.series[0]
            for i, point in enumerate(series.points):
                fill = point.format.fill
                fill.solid()
                fill.fore_color.rgb = self.SEGMENT_ORANGE

        # Insight
        insight = data.get('insight', '')
        self._add_insight_box(slide, insight)

    def _add_team_slide(self):
        """Slide 9: Team Activity"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        data = self.data.get('slide_9_team', {})
        self._add_slide_title(slide, data.get('title', 'Team: Active Workspace Users'))

        charts = data.get('charts', {})

        # Active Users Count
        active_users = charts.get('active_users_last_90_days', {})
        user_count = active_users.get('count', 0)

        user_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(2), Inches(1))
        uf = user_box.text_frame
        uf.text = f"{user_count}\nActive Users"
        uf.paragraphs[0].font.size = Pt(36)
        uf.paragraphs[0].font.bold = True
        uf.paragraphs[0].font.color.rgb = self.SEGMENT_PURPLE
        uf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Workspace Events
        events = charts.get('workspace_events', {})
        y_pos = 1.8

        for event_type, count in events.items():
            label = event_type.replace('_', ' ').title()
            event_box = slide.shapes.add_textbox(Inches(3.5), Inches(y_pos), Inches(5), Inches(0.4))
            ef = event_box.text_frame
            ef.text = f"• {label}: {count}"
            ef.paragraphs[0].font.size = Pt(14)
            y_pos += 0.4

        # Insight
        insight = data.get('insight', '')
        self._add_insight_box(slide, insight)

    def _add_summary_slide(self):
        """Slide 11: Summary"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        data = self.data.get('slide_11_summary', {})
        self._add_slide_title(slide, data.get('title', 'Summary'))

        # Health Ratings
        ratings = data.get('health_ratings', {})

        y_pos = 1.8
        for metric, rating in ratings.items():
            label = metric.replace('_', ' ').title()

            rating_box = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.5))
            rf = rating_box.text_frame
            rf.text = f"{label}: {rating}"
            rf.paragraphs[0].font.size = Pt(18)

            # Color based on rating
            if '✓' in rating:
                rf.paragraphs[0].font.color.rgb = self.SEGMENT_GREEN
            elif '⚠' in rating:
                rf.paragraphs[0].font.color.rgb = self.SEGMENT_YELLOW
            elif '✗' in rating:
                rf.paragraphs[0].font.color.rgb = self.SEGMENT_RED

            y_pos += 0.6

    def _add_conclusions_slide(self):
        """Slide 12: Conclusions"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        data = self.data.get('slide_12_conclusions', {})
        self._add_slide_title(slide, data.get('title', 'Conclusions'))

        # Three columns
        columns = [
            ('Biggest Issues', data.get('biggest_issues', [])),
            ('Possible Causes', data.get('possible_causes', [])),
            ('How to Address', data.get('how_to_address', []))
        ]

        x_positions = [Inches(0.5), Inches(3.5), Inches(6.5)]

        for idx, (title, items) in enumerate(columns):
            # Column title
            title_box = slide.shapes.add_textbox(x_positions[idx], Inches(1.5), Inches(2.8), Inches(0.4))
            tf = title_box.text_frame
            tf.text = title
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = self.SEGMENT_BLUE

            # Items
            items_box = slide.shapes.add_textbox(x_positions[idx], Inches(2), Inches(2.8), Inches(3))
            if_frame = items_box.text_frame
            if_frame.word_wrap = True

            for item in items[:3]:  # Top 3
                p = if_frame.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(11)
                p.space_after = Pt(8)

    def _add_thank_you_slide(self):
        """Slide 13: Thank You"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Thank you text
        thanks = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1))
        tf = thanks.text_frame
        tf.text = "Thank You!\nLet's discuss your questions."
        tf.paragraphs[0].font.size = Pt(44)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_slide_title(self, slide, title: str):
        """Add title to slide"""
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        tf = title_box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(33, 33, 33)

    def _add_insight_box(self, slide, insight: str, y_position: float = 4.8):
        """Add insight text box at bottom of slide"""
        if insight:
            insight_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_position), Inches(9), Inches(0.7))
            if_frame = insight_box.text_frame
            if_frame.text = f"💡 {insight}"
            if_frame.word_wrap = True
            if_frame.paragraphs[0].font.size = Pt(12)
            if_frame.paragraphs[0].font.italic = True
            if_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)


def main():
    """CLI entry point"""
    if len(sys.argv) < 2:
        print("Usage: python generate_health_check_ppt.py <gem_output.json> [customer_name]")
        print("\nExample:")
        print("  python generate_health_check_ppt.py gem_output.json 'Mission Lane'")
        sys.exit(1)

    json_path = sys.argv[1]
    customer_name = sys.argv[2] if len(sys.argv) > 2 else "Customer"

    if not Path(json_path).exists():
        print(f"❌ Error: File not found: {json_path}")
        sys.exit(1)

    try:
        generator = HealthCheckPPTGenerator(json_path, customer_name)
        output_path = generator.generate()
        print(f"\n✅ Success! Open: {output_path}")
    except Exception as e:
        print(f"❌ Error generating PowerPoint: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == '__main__':
    main()
